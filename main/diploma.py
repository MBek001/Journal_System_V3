import qrcode
from datetime import datetime

from PIL import Image
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx import Presentation
from pptx.util import Pt, Cm
import os

from JournalSystem import settings
from JournalSystem.settings import MEDIA_ROOT


def create_diploma(author_name, journal_name, issue_number, article_url, pub_date, template_path, output_path):
    prs = Presentation(template_path)
    slide = prs.slides[0]

    # --- Author Name ---
    left = Pt(150)
    top = Pt(175)
    width = Pt(600)
    height = Pt(50)

    author_box = slide.shapes.add_textbox(left, top, width, height)
    author_tf = author_box.text_frame
    author_tf.clear()
    p = author_tf.add_paragraph()
    p.text = author_name.upper()
    p.font.size = Pt(32)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # --- Article Info ---
    article_info = f"«{journal_name}» ЖУРНАЛИНИНГ\n{issue_number}-СОНИДА ИЛМИЙ МАҚОЛАСИ ЧОП ЭТИЛГАНЛИГИ УЧУН"

    left2 = Pt(100)
    top2 = Pt(235)
    width2 = Pt(650)
    height2 = Pt(80)

    info_box = slide.shapes.add_textbox(left2, top2, width2, height2)
    info_tf = info_box.text_frame
    info_tf.clear()
    p2 = info_tf.add_paragraph()
    p2.text = article_info
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(0, 51, 153)
    p2.alignment = PP_ALIGN.CENTER

    if article_url:
        qr_dir = os.path.join(settings.MEDIA_ROOT, "qr_codes")
        os.makedirs(qr_dir, exist_ok=True)

        qr_path = os.path.join(qr_dir, "_qr.png")

        # Generate QR
        qr_img = qrcode.make(article_url)
        qr_img.save(qr_path)
        Image.open(qr_path).convert("RGB").save(qr_path, "PNG")

        if os.path.exists(qr_path):
            print("Adding QR from:", qr_path)
            qr_left = Cm(21)
            qr_top = Cm(15)
            qr_size = Cm(4)

            slide.shapes.add_picture(qr_path, qr_left, qr_top, qr_size, qr_size)

    # --- Date under QR code ---
    if not pub_date:
        pub_date = datetime.today()

    if isinstance(pub_date, str):
        try:
            pub_date = datetime.strptime(pub_date, "%Y-%m-%d")
        except:
            pub_date = datetime.today()

    date_text = pub_date.strftime("%d.%m.%Y йил.")
    date_left = Cm(20.5)
    date_top = Cm(17.8)
    date_width = Cm(5)
    date_height = Cm(1)

    date_box = slide.shapes.add_textbox(date_left, date_top, date_width, date_height)
    date_tf = date_box.text_frame
    date_tf.clear()
    p3 = date_tf.add_paragraph()
    p3.text = date_text
    p3.font.size = Pt(12)
    p3.alignment = PP_ALIGN.CENTER

    prs.save(output_path)
    return output_path
